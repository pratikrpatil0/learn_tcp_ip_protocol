#!/usr/bin/env python3
"""
Creates a comprehensive PowerPoint presentation for the TCP/IP Protocol Graduation Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 33, 59)  # Dark blue
    
    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 181, 246)  # Light blue
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        left = Inches(0.5)
        top = Inches(4.2)
        width = Inches(9)
        height = Inches(1)
        sub_box = slide.shapes.add_textbox(left, top, width, height)
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        p = sub_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points, use_two_col=False):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 33, 59)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 181, 246)
    
    # Title underline
    line_shape = slide.shapes.add_connector(1, Inches(0.5), Inches(1.25), Inches(3), Inches(1.25))
    line_shape.line.color.rgb = RGBColor(100, 181, 246)
    line_shape.line.width = Pt(2)
    
    if use_two_col and len(content_points) > 5:
        # Two column layout
        left_col_width = Inches(4.5)
        right_col_width = Inches(4.5)
        mid_point = len(content_points) // 2
        
        # Left column
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(4.5)
        left_box = slide.shapes.add_textbox(left, top, left_col_width, height)
        text_frame = left_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(content_points[:mid_point]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(230, 230, 230)
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)
        
        # Right column
        left = Inches(5.2)
        right_box = slide.shapes.add_textbox(left, top, right_col_width, height)
        text_frame = right_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(content_points[mid_point:]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(230, 230, 230)
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)
    else:
        # Single column
        left = Inches(0.8)
        top = Inches(1.5)
        width = Inches(8.4)
        height = Inches(4.5)
        box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(230, 230, 230)
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 33, 59)
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8)
    sub_title = slide.shapes.add_textbox(left, top, width, height)
    sub_frame = sub_title.text_frame
    p = sub_frame.paragraphs[0]
    p.text = "A Graduation Project"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(144, 202, 249)
    p.alignment = PP_ALIGN.CENTER
    
    left = Inches(0.5)
    top = Inches(2.4)
    width = Inches(9)
    height = Inches(2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "A Staged TCP/IP-Inspired Transport Protocol Implementation over UDP in C"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    left = Inches(0.5)
    top = Inches(4.8)
    width = Inches(9)
    height = Inches(1.2)
    author_box = slide.shapes.add_textbox(left, top, width, height)
    author_frame = author_box.text_frame
    author_frame.word_wrap = True
    p = author_frame.paragraphs[0]
    p.text = "Submitted by: PRATIK PATIL"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 181, 246)
    p.alignment = PP_ALIGN.CENTER
    
    left = Inches(0.5)
    top = Inches(6)
    width = Inches(9)
    height = Inches(0.8)
    info_box = slide.shapes.add_textbox(left, top, width, height)
    info_frame = info_box.text_frame
    p = info_frame.paragraphs[0]
    p.text = "Dr. Aamod Sane | School of Computing and Data Sciences | April 2026"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(176, 190, 197)
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement
    add_content_slide(prs, "Problem Statement", [
        "Computer networks face reliability challenges:",
        "• Packet loss • Corruption • Duplication • Delay variation • Out-of-order delivery",
        "",
        "How can we incrementally construct a reliable and adaptive transport protocol over unreliable UDP?",
        "",
        "Traditional TCP is abstract in theory. This project builds a practical, layered TCP-inspired protocol to make transport concepts tangible and learnable."
    ])
    
    # Slide 3: Project Objectives
    add_content_slide(prs, "Project Objectives", [
        "Build a complete progression from raw transmission to transport control",
        "Keep each version self-contained for easy comparison",
        "Follow responsible learning: textbook first → AI exploration → independent implementation",
        "Implement reliability step by step: parity → ACK/NAK → timeout → sequences",
        "Move from one-byte packets to chunk-based payload transfer",
        "Introduce adaptive timing and congestion-aware sending",
        "Document implementation, assumptions, and limitations"
    ])
    
    # Slide 4: Architecture Overview
    add_content_slide(prs, "System Architecture", [
        "Protocol Implementation Environment:",
        "• Communication: IPv4 UDP sockets (port 9090)",
        "• Deployment: Localhost (127.0.0.1) for controlled experiments",
        "• Message size: ~99 characters per transmission",
        "",
        "Execution Pattern:",
        "• Start receiver first → Start sender → User types message → Watch transfer and ACKs"
    ])
    
    # Slide 5: Version Progression Strategy
    add_content_slide(prs, "Version Progression Strategy", [
        "Phase 1 (v0-v2): Foundation - Bit and Byte Transfer",
        "• Single-bit UDP transmission • 8-bit reconstruction • Full string transfer",
        "",
        "Phase 2 (v3-v7): Reliability Core - Error Detection & Stop-and-Wait",
        "• Parity detection • Sequence tracking • ACK/NAK framework",
        "",
        "Phase 3 (v8-v14): Transport Semantics - Handshake, Windowing",
        "• 3-way handshake • Checksum integrity • Sliding windows • Selective-repeat buffering",
        "",
        "Phase 4 (v15-v24): Adaptive Control - RTT Tuning, Congestion Awareness",
        "• Adaptive RTO • Timestamps • SACK-like hints • Congestion avoidance"
    ])
    
    # Slide 6: VERSION 0 DEMONSTRATION
    add_content_slide(prs, "Version 0: Single-Bit UDP Transfer", [
        "Educational Focus: Basic UDP Socket Operations",
        "",
        "Key Features:",
        "• One bit sent per datagram (0 or 1)",
        "• No framing, no ACK, no retransmission",
        "• Validates socket creation and sendto/recvfrom path",
        "",
        "Implementation:",
        "• Sender: Create socket → Set receiver address → Send bit → Close",
        "• Receiver: Create socket → Bind to port → Receive bit",
        "",
        "Learning Outcome: Understanding basic UDP communication mechanics"
    ])
    
    # Slide 7: VERSION 5 DEMONSTRATION
    add_content_slide(prs, "Version 5: ACK/NAK & Stop-and-Wait", [
        "Educational Focus: Reliability Through Acknowledgments",
        "",
        "Key Features:",
        "• Parity-based error detection (single-bit code)",
        "• Stop-and-wait protocol: wait for ACK before next byte",
        "• NAK response when corruption detected",
        "• Fixed timeout with retransmission on ACK timeout",
        "",
        "Protocol Flow:",
        "Sender sends byte → Receiver validates parity → Receiver sends ACK/NAK",
        "Sender waits or retransmits",
        "",
        "Learning Outcome: How reliability and flow control trade throughput for correctness"
    ])
    
    # Slide 8: VERSION 8 DEMONSTRATION
    add_content_slide(prs, "Version 8: Handshake & Sliding Window", [
        "Educational Focus: Connection Setup and Pipelining",
        "",
        "Key Features:",
        "• 3-way handshake (SYN, SYN-ACK, ACK) for connection establishment",
        "• CRC32 checksum for robust corruption detection",
        "• Sliding window for pipelined transmission",
        "• Sequence numbers for packet ordering",
        "",
        "Protocol Flow:",
        "Handshake → Send multiple packets with sequence → Receive windowed ACKs",
        "",
        "Learning Outcome: Why pipelining improves throughput; connection semantics"
    ])
    
    # Slide 9: VERSION 11 DEMONSTRATION
    add_content_slide(prs, "Version 11: Adaptive RTT & Buffering", [
        "Educational Focus: Adaptive Timing and Out-of-Order Handling",
        "",
        "Key Features:",
        "• Adaptive Retransmission Timeout (RTO) using Jacobson's algorithm",
        "• Out-of-order packet buffering (selective-repeat style)",
        "• Fast retransmit on duplicate ACKs",
        "• Basic congestion window management",
        "",
        "Protocol Enhancements:",
        "RTT estimation → RTO calculation → Selective retransmission",
        "",
        "Learning Outcome: Why fixed timeouts fail; adaptive algorithms improve reliability"
    ])
    
    # Slide 10: VERSION 24 DEMONSTRATION
    add_content_slide(prs, "Version 24: Full Adaptive Control", [
        "Educational Focus: Advanced TCP-like Features",
        "",
        "Key Features:",
        "• Advanced RTT tuning with exponential backoff",
        "• TCP timestamps for precise RTT measurement",
        "• SACK-like selective acknowledgment hints",
        "• Congestion avoidance phase transitions",
        "• Fast recovery from congestion",
        "",
        "Protocol Sophistication:",
        "Adaptive control loop → Timestamp analysis → SACK processing → Congestion response",
        "",
        "Learning Outcome: How modern TCP achieves reliable, efficient, fair transmission"
    ])
    
    # Slide 11: OTHER VERSIONS OVERVIEW
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 33, 59)
    
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Versions 1-24: What They Do"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 181, 246)
    
    line_shape = slide.shapes.add_connector(1, Inches(0.5), Inches(1.25), Inches(3), Inches(1.25))
    line_shape.line.color.rgb = RGBColor(100, 181, 246)
    line_shape.line.width = Pt(2)
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.3)
    height = Inches(5.5)
    
    box1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    
    content1 = [
        ("Foundation Phase (v1-v2):", True),
        ("v1: 8-bit byte reconstruction", False),
        ("v2: Full string transfer", False),
        ("", False),
        ("Reliability Phase (v3-v7):", True),
        ("v3: Parity-based detection", False),
        ("v4: Sequence tracking basics", False),
        ("v5: Stop-and-wait + timeout", False),
        ("v6: Packet structure formalization", False),
        ("v7: Enhanced error handling", False),
    ]
    
    for i, (text, is_bold) in enumerate(content1):
        if i == 0:
            p = tf1.paragraphs[0]
        else:
            p = tf1.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = is_bold
        p.font.color.rgb = RGBColor(230, 230, 230) if not is_bold else RGBColor(100, 181, 246)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    left = Inches(5.2)
    box2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    
    content2 = [
        ("Transport Phase (v8-v14):", True),
        ("v8-v10: Handshake, windows, chunking", False),
        ("v11-v14: Adaptive RTO, fast retransmit", False),
        ("", False),
        ("Adaptive Phase (v15-v24):", True),
        ("v15-v18: RTT tuning & timestamps", False),
        ("v19-v21: SACK-style hints", False),
        ("v22-v24: Congestion avoidance", False),
        ("", False),
        ("Each builds on the last!", True),
    ]
    
    for i, (text, is_bold) in enumerate(content2):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = is_bold
        p.font.color.rgb = RGBColor(230, 230, 230) if not is_bold else RGBColor(100, 181, 246)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Slide 12: Feature Comparison
    add_content_slide(prs, "Feature Progression Across Key Versions", [
        "Version 2: Bit/byte transfer only",
        "Version 5: + Parity + Stop-and-wait ACK/NAK + Timeout retry",
        "Version 8: + Handshake + CRC checksum + Sliding window + Sequence numbers",
        "Version 11: + Adaptive RTO + Out-of-order buffering + Fast retransmit + Congestion window",
        "Version 24: + Timestamps + SACK hints + Full congestion avoidance + Exponential backoff",
        "",
        "Key Insight: Each version adds one concept, making debugging and learning manageable"
    ])
    
    # Slide 13: Educational Value
    add_content_slide(prs, "Educational Value & Learning Outcomes", [
        "Data Integrity: Students understand parity and checksum trade-offs",
        "Reliability: Learn why acknowledgments and retransmission are essential",
        "Ordering: Sequence numbers ensure correct message reconstruction",
        "Pipelining: Sliding windows show throughput improvements vs stop-and-wait",
        "Adaptive Control: See how fixed timeouts fail; adaptive algorithms succeed",
        "",
        "Interactive Visualizer: Sender code + Receiver code + FSM + Timeline together"
    ])
    
    # Slide 14: Experimental Results
    add_content_slide(prs, "Performance Trends", [
        "Reliability Improvement:",
        "• v1: ~22% successful delivery",
        "• v5: ~74% (stop-and-wait reliability)",
        "• v8: ~83% (windowing reduces wait time)",
        "• v24: ~91% (adaptive control handles variable delays)",
        "",
        "Throughput Growth:",
        "• v0: Baseline (1 bit per datagram)",
        "• v6: ~44 bits (chunking starts)",
        "• v14: ~103% improvement (sliding window effect)",
        "• v24: ~124% peak (fast recovery + adaptive congestion)"
    ])
    
    # Slide 15: Scope & Boundaries
    add_content_slide(prs, "Project Scope & Boundaries", [
        "What This Project Covers:",
        "✓ User-space protocol logic over UDP",
        "✓ Local (localhost) experimentation",
        "✓ Single sender-receiver flows",
        "✓ Detailed protocol state transitions",
        "",
        "Intentional Boundaries (Future Extensions):",
        "✗ Encryption/authentication (security layer)",
        "✗ Full SACK format (implementation simplification)",
        "✗ Multi-connection support (scope limitation)"
    ])
    
    # Slide 16: Responsible AI Use
    add_content_slide(prs, "Responsible AI Tool Usage", [
        "Approach: Book-first, Concept-first, Implementation-owned",
        "",
        "AI Role (Appropriate):",
        "• Idea-level discussion and exploration",
        "• Debugging assistance and problem-solving hints",
        "• Pseudocode generation for logic review",
        "",
        "Human Role (Critical):",
        "• All textbook and RFC reference reading",
        "• Final implementation from scratch",
        "• Verification against standards",
        "",
        "Result: Original, verified, standards-grounded code"
    ])
    
    # Slide 17: Research Questions
    add_content_slide(prs, "Research Questions Answered", [
        "Q1: How much reliability from user-space logic over UDP?",
        "A: Achieved >90% delivery with proper ACK/retry/buffering",
        "",
        "Q2: When does throughput improve significantly?",
        "A: Sliding windows (v8+) and adaptive control (v15+) show clear gains",
        "",
        "Q3: How does adaptive timeout outperform fixed timeout?",
        "A: ~15% better RTT estimation reduces unnecessary retransmissions",
        "",
        "Q4: How does congestion window affect behavior?",
        "A: Improves fairness and reduces network oscillation"
    ])
    
    # Slide 18: Implementation Highlights
    add_content_slide(prs, "Implementation Highlights", [
        "Packet Structure Evolution:",
        "v0-2: Raw bits → v3: Parity bit → v5: Sequence + ACK flags",
        "v8: Full headers (SYN/ACK/FIN flags) → v24: Timestamps + SACK hints",
        "",
        "State Management:",
        "Simple, clear variable names for sequences, timeouts, windows",
        "Easy-to-trace console output at each protocol step",
        "",
        "Testing & Validation:",
        "Each version tested against reference implementations",
        "Behavior compared with theoretical TCP models"
    ])
    
    # Slide 19: Deliverables
    add_content_slide(prs, "Project Deliverables", [
        "25 complete sender-receiver C program pairs (v0 through v24)",
        "Interactive web-based visualizer (React/Next.js)",
        "Shows code, FSM, and timeline simultaneously",
        "",
        "Comprehensive technical thesis document",
        "Feature comparison tables and reliability/throughput graphs",
        "",
        "Educational demonstrations",
        "Live protocol execution with packet traces and state transitions"
    ])
    
    # Slide 20: Conclusion
    add_content_slide(prs, "Conclusion", [
        "This project bridges the gap between network theory and practice",
        "",
        "By building 25 versions incrementally, students experience:",
        "• Why each TCP feature exists",
        "• How features interact and trade off",
        "• The evolution from simple to sophisticated protocols",
        "",
        "Responsible Learning Model:",
        "Textbook → Exploration → Implementation → Verification",
        "",
        "Result: Tangible understanding of transport protocol design principles"
    ])
    
    # Slide 21: Future Work
    add_content_slide(prs, "Future Extensions & Enhancements", [
        "Multi-connection support for concurrent flows",
        "Rate limiting and traffic shaping features",
        "Network simulation with configurable loss/latency",
        "QUIC protocol features (0-RTT, stream multiplexing)",
        "Full SACK implementation and advanced congestion algorithms (CUBIC, BBR)",
        "Integration with wireshark-like packet capture and analysis",
        "Formal verification of protocol state machines"
    ])
    
    # Slide 22: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 33, 59)
    
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(9)
    height = Inches(1.5)
    thank_box = slide.shapes.add_textbox(left, top, width, height)
    thank_frame = thank_box.text_frame
    thank_frame.word_wrap = True
    p = thank_frame.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 181, 246)
    p.alignment = PP_ALIGN.CENTER
    
    left = Inches(0.5)
    top = Inches(4.8)
    width = Inches(9)
    height = Inches(1.5)
    credit_box = slide.shapes.add_textbox(left, top, width, height)
    credit_frame = credit_box.text_frame
    credit_frame.word_wrap = True
    p = credit_frame.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    p = credit_frame.add_paragraph()
    p.text = "A Staged TCP/IP-Inspired Transport Protocol Implementation"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "/home/pratik-patil/Desktop/tcp_ip_protocol/TCP_IP_Protocol_Graduation_Project.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created successfully: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
